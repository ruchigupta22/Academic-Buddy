import io
from pathlib import Path
from typing import List, Dict, Any
 
import fitz  # PyMuPDF — the 'fitz' name is historical
from pptx import Presentation
 
 
def parse_pdf(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    
    pages = []
 
    # Open PDF from bytes (no temp file needed)
    doc = fitz.open(stream=file_bytes, filetype="pdf")
 
    for page_num, page in enumerate(doc, start=1):
        # get_text("text") returns plain text, preserving paragraph breaks
        # get_text("blocks") would give bounding boxes too — useful for tables
        text = page.get_text("text").strip()
 
        # Skip pages that are blank or image-only (no extractable text)
        if len(text) < 30:
            continue
 
        pages.append({
            "text": text,
            "page": page_num,
            "source": filename,
        })
 
    doc.close()
    return pages
 
 
def parse_pptx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    
    slides = []
 
    prs = Presentation(io.BytesIO(file_bytes))
 
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
 
        # Extract text from every shape on the slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs).strip()
                if line:
                    texts.append(line)
 
        # Also grab speaker notes (often the richest content)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = notes_frame.text.strip()
            if notes and notes != "Click to edit Master text styles":
                texts.append(f"[Speaker notes]: {notes}")
 
        combined = "\n".join(texts).strip()
        if len(combined) < 20:
            continue
 
        slides.append({
            "text": combined,
            "page": slide_num,
            "source": filename,
        })
 
    return slides
 
 
def parse_file(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Router function — dispatches to the right parser based on file extension.
    This is the only function other modules call; they don't know the file type.
    """
    ext = Path(filename).suffix.lower()
 
    if ext == ".pdf":
        return parse_pdf(file_bytes, filename)
    elif ext in (".pptx", ".ppt"):
        return parse_pptx(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Upload PDF or PPTX.")
 